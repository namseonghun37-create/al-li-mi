"""Generate the TribeLink startup pitch deck (.pptx)."""

import os

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import Circle

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt

# --------------------------------------------------------------------------- #
# Design system
# --------------------------------------------------------------------------- #
NAVY = RGBColor(0x0D, 0x1B, 0x2A)        # background
CARD = RGBColor(0x1E, 0x2D, 0x3D)        # card / box background
TEAL = RGBColor(0x00, 0xC9, 0xA7)        # accent 1
GOLD = RGBColor(0xFF, 0xD7, 0x00)        # accent 2 (CTA only)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)       # primary text
GREY = RGBColor(0xB0, 0xBE, 0xC5)        # secondary text
WARM = RGBColor(0xFF, 0x6B, 0x4A)        # problem accent (red/orange)
DARK_TEXT = RGBColor(0x0D, 0x1B, 0x2A)   # text on white cards
CARD_DESC = RGBColor(0x5A, 0x6B, 0x7B)   # description on white cards
GREEN_OK = RGBColor(0x6F, 0xCF, 0x97)    # competitor checkmark
MUTED = RGBColor(0x6B, 0x7A, 0x88)       # competitor cross

FONT = "Montserrat"          # primary (PowerPoint falls back to Calibri)
FONT_FB = "Calibri"

SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

ASSETS = "assets"
OUTPUT = "/mnt/user-data/outputs/TribeLink_Pitch_Deck.pptx"

# matplotlib defaults for transparent dark-slide charts
plt.rcParams.update({
    "text.color": "white",
    "axes.edgecolor": "white",
    "font.family": "DejaVu Sans",
    "svg.fonttype": "none",
})


# --------------------------------------------------------------------------- #
# Low-level helpers
# --------------------------------------------------------------------------- #
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None,
             shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except (IndexError, KeyError):
            pass
    return sp


def add_shadow(shape, blur=Pt(12), dist=Pt(4), alpha=58000, color="000000"):
    """Inject a soft outer shadow into an autoshape."""
    spPr = shape._element.spPr
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    shdw = eff.makeelement(qn("a:outerShdw"), {
        "blurRad": str(int(blur)),
        "dist": str(int(dist)),
        "dir": "5400000",
        "rotWithShape": "0",
    })
    clr = shdw.makeelement(qn("a:srgbClr"), {"val": color})
    a = clr.makeelement(qn("a:alpha"), {"val": str(alpha)})
    clr.append(a)
    shdw.append(clr)
    eff.append(shdw)
    spPr.append(eff)


def text_box(slide, x, y, w, h, runs, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT,
             line_spacing=None, italic=False, wrap=True):
    """runs may be a string or a list of paragraphs.

    Each paragraph is either a string or a list of (text, overrides) run tuples,
    where overrides is a dict that can set size/bold/color/italic/font.
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    if isinstance(runs, str):
        runs = [runs]

    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        chunks = [para] if isinstance(para, str) else para
        for chunk in chunks:
            if isinstance(chunk, str):
                txt, ov = chunk, {}
            else:
                txt, ov = chunk
            r = p.add_run()
            r.text = txt
            f = r.font
            f.size = Pt(ov.get("size", size))
            f.bold = ov.get("bold", bold)
            f.italic = ov.get("italic", italic)
            f.name = ov.get("font", font)
            f.color.rgb = ov.get("color", color)
    return tb


def base_slide(prs, page_num, bg=NAVY, border=TEAL):
    """Create a slide with background, left accent border and slide number."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=bg)            # background
    add_rect(slide, 0, 0, Cm(0.4), SLIDE_H, fill=border)        # left border
    # slide number bottom-right
    text_box(slide, Cm(31.0), Cm(18.1), Cm(2.5), Cm(0.8),
             f"{page_num:02d} / 10", size=10, color=GREY,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def headline(slide, text, size=30, color=WHITE, top=1.05):
    add_rect(slide, Cm(1.6), Cm(top + 0.05), Cm(0.5), Cm(1.0), fill=TEAL)
    text_box(slide, Cm(2.4), Cm(top), Cm(30), Cm(1.5), text,
             size=size, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)


def style_cell(cell, text, size=14, bold=False, color=WHITE, fill=CARD,
               align=PP_ALIGN.CENTER):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Cm(0.15)
    cell.margin_right = Cm(0.15)
    cell.margin_top = Cm(0.05)
    cell.margin_bottom = Cm(0.05)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = FONT
    r.font.color.rgb = color


# --------------------------------------------------------------------------- #
# Matplotlib chart assets
# --------------------------------------------------------------------------- #
def make_overlap_circles(path):
    fig, ax = plt.subplots(figsize=(9, 7))
    ax.set_xlim(-2.2, 2.2)
    ax.set_ylim(-2.0, 2.2)
    ax.set_aspect("equal")
    ax.axis("off")

    specs = [
        ((0.0, 0.55), "#00C9A7", "AI Matching", (0.0, 1.95)),
        ((-0.62, -0.45), "#1E90B8", "Virtual HQ", (-1.55, -1.45)),
        ((0.62, -0.45), "#6C63FF", "Tribe Community", (1.55, -1.45)),
    ]
    r = 0.95
    for (cx, cy), col, _, _ in specs:
        ax.add_patch(Circle((cx, cy), r, facecolor=col, edgecolor="white",
                            linewidth=2, alpha=0.55))
    for (cx, cy), col, label, (lx, ly) in specs:
        ax.text(lx, ly, label, ha="center", va="center",
                fontsize=20, fontweight="bold", color="white")
    ax.text(0.0, -0.12, "TribeLink", ha="center", va="center",
            fontsize=16, fontweight="bold", color="white", alpha=0.95)
    fig.savefig(path, transparent=True, bbox_inches="tight", dpi=200)
    plt.close(fig)


def make_concentric(path):
    fig, ax = plt.subplots(figsize=(8, 8))
    ax.set_xlim(-1.15, 1.15)
    ax.set_ylim(-1.15, 1.15)
    ax.set_aspect("equal")
    ax.axis("off")

    rings = [
        (1.0, "#0E3A40", "TAM", "$45B", "Global remote\nwork software", 0.78),
        (0.66, "#11857A", "SAM", "$12B", "SMEs & freelance\ncollectives", 0.45),
        (0.34, "#00C9A7", "SOM", "$480M", "Initial target\n(Years 1–3)", 0.0),
    ]
    for r, col, *_ in rings:
        ax.add_patch(Circle((0, 0), r, facecolor=col, edgecolor="white",
                            linewidth=1.5, zorder=1))
    for r, col, lbl, val, desc, ty in rings:
        if lbl == "SOM":
            ax.text(0, 0.07, lbl, ha="center", va="center", fontsize=15,
                    fontweight="bold", color="white")
            ax.text(0, -0.07, val, ha="center", va="center", fontsize=14,
                    fontweight="bold", color="white")
        else:
            ax.text(0, ty + 0.06, f"{lbl}  {val}", ha="center", va="center",
                    fontsize=15, fontweight="bold", color="white")
            ax.text(0, ty - 0.07, desc, ha="center", va="center", fontsize=9.5,
                    color="white")
    fig.savefig(path, transparent=True, bbox_inches="tight", dpi=200)
    plt.close(fig)


def make_donut(path):
    fig, ax = plt.subplots(figsize=(7.4, 7.4))
    sizes = [40, 30, 20, 10]
    labels = ["Product\nDevelopment", "Sales &\nMarketing",
              "Operations\n& Talent", "Legal &\nCompliance"]
    colors = ["#2D6CDF", "#00C9A7", "#3CB371", "#90A4AE"]
    wedges, _ = ax.pie(sizes, colors=colors, startangle=90,
                       counterclock=False,
                       wedgeprops=dict(width=0.36, edgecolor="#0D1B2A",
                                       linewidth=3))
    # percentage labels on the ring
    import numpy as np
    for w, pct in zip(wedges, sizes):
        ang = np.deg2rad((w.theta2 + w.theta1) / 2.0)
        ax.text(0.82 * np.cos(ang), 0.82 * np.sin(ang), f"{pct}%",
                ha="center", va="center", fontsize=14, fontweight="bold",
                color="white")
    ax.text(0, 0, "$500K", ha="center", va="center", fontsize=24,
            fontweight="bold", color="white")
    ax.set_aspect("equal")
    fig.savefig(path, transparent=True, bbox_inches="tight", dpi=200)
    plt.close(fig)


# --------------------------------------------------------------------------- #
# Slides
# --------------------------------------------------------------------------- #
def slide_title(prs):
    s = base_slide(prs, 1)
    # logo placeholder circle top-right
    logo = add_rect(s, Cm(30.4), Cm(1.1), Cm(2.4), Cm(2.4),
                    fill=None, line=TEAL, line_w=Pt(2.25), shape=MSO_SHAPE.OVAL)
    text_box(s, Cm(30.4), Cm(1.1), Cm(2.4), Cm(2.4), "LOGO", size=10,
             color=TEAL, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    text_box(s, Cm(2.5), Cm(6.4), Cm(28.8), Cm(2.8), "TribeLink",
             size=66, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Cm(14.43), Cm(9.55), Cm(5.0), Cm(0.09), fill=TEAL)
    text_box(s, Cm(2.5), Cm(9.9), Cm(28.8), Cm(1.4),
             "“Where Global Teams Become Tribes”", size=24,
             color=TEAL, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, Cm(2.5), Cm(16.4), Cm(28.8), Cm(1.0),
             "Presented by  Speaker A (CEO)   |   Speaker B (CPO)   |   "
             "Speaker C (CFO)", size=14, color=GREY, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


def slide_problem(prs):
    s = base_slide(prs, 2)
    headline(s, "Remote Workers Are Drowning", size=32)

    cards = [
        ("\U0001F527", "Tool Overload", "Average of 8 disconnected\ntools per worker"),
        ("\U0001F494", "Disconnection", "87% of remote workers\nfeel isolated"),
        ("\U0001F4B8", "Talent Loss", "Billions lost every year\nto disengagement"),
    ]
    n = 3
    gap = Cm(0.9)
    left = Cm(1.6)
    total = SLIDE_W - left - Cm(1.2)
    cw = (total - gap * (n - 1)) / n
    top = Cm(5.4)
    ch = Cm(8.4)
    for i, (icon, title, desc) in enumerate(cards):
        x = left + i * (cw + gap)
        card = add_rect(s, x, top, cw, ch, fill=CARD,
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
        add_shadow(card)
        text_box(s, x, top + Cm(0.9), cw, Cm(2.2), icon, size=58,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x, top + Cm(3.5), cw, Cm(1.2), title, size=22, bold=True,
                 color=WARM, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x + Cm(0.5), top + Cm(4.9), cw - Cm(1.0), Cm(2.6), desc,
                 size=15, color=GREY, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.TOP, line_spacing=1.1)

    text_box(s, Cm(1.6), Cm(17.1), Cm(30), Cm(0.9),
             "Based on a survey of 1,200+ remote professionals across 30 "
             "countries.", size=11, color=GREY, italic=True)


def slide_solution(prs, img):
    s = base_slide(prs, 3)
    headline(s, "Introducing TribeLink", size=32)
    iw = Cm(18.5)
    ih = iw * 7 / 9
    s.shapes.add_picture(img, (SLIDE_W - iw) / 2, Cm(3.7), width=iw, height=ih)
    text_box(s, Cm(2.5), Cm(16.7), Cm(28.8), Cm(1.2),
             "One platform. One tribe. Infinite potential.", size=22,
             bold=True, color=TEAL, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


def slide_features(prs):
    s = base_slide(prs, 4)
    headline(s, "4 Features That Set Us Apart", size=32)

    cards = [
        ("\U0001F916", "AI Tribe Matching", "Accurate, adaptive role pairing"),
        ("\U0001F5F3️", "Smart Voting System", "Democratic team decisions"),
        ("\U0001F4CA", "Contribution Analytics", "Real-time impact tracking"),
        ("\U0001F33F", "Green Workspace",
         "Environmentally friendly carbon dashboard"),
    ]
    left = Cm(1.6)
    top = Cm(5.4)
    gx, gy = Cm(0.9), Cm(0.8)
    total = SLIDE_W - left - Cm(1.2)
    cw = (total - gx) / 2
    ch = Cm(5.4)
    for i, (icon, title, desc) in enumerate(cards):
        col, row = i % 2, i // 2
        x = left + col * (cw + gx)
        y = top + row * (ch + gy)
        card = add_rect(s, x, y, cw, ch, fill=WHITE,
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
        add_shadow(card, alpha=40000)
        add_rect(s, x, y, Cm(0.3), ch, fill=TEAL)  # left accent strip
        text_box(s, x + Cm(0.9), y + Cm(0.7), Cm(3.0), Cm(2.0), icon, size=40,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x + Cm(4.2), y + Cm(1.0), cw - Cm(4.7), Cm(1.4), title,
                 size=21, bold=True, color=DARK_TEXT, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x + Cm(4.2), y + Cm(2.6), cw - Cm(4.7), Cm(2.2), desc,
                 size=14, color=CARD_DESC, anchor=MSO_ANCHOR.TOP,
                 line_spacing=1.1)


def slide_market(prs, img):
    s = base_slide(prs, 5)
    headline(s, "A $45B Global Opportunity", size=32)

    ih = Cm(12.0)
    s.shapes.add_picture(img, Cm(2.0), Cm(5.0), height=ih)

    # right stat callouts
    stats = [("1.57B", "freelancers worldwide"),
             ("15%", "CAGR through 2030")]
    sx = Cm(19.5)
    sw = Cm(12.4)
    sh = Cm(4.6)
    sy0 = Cm(5.6)
    for i, (num, lbl) in enumerate(stats):
        y = sy0 + i * (sh + Cm(1.1))
        card = add_rect(s, sx, y, sw, sh, fill=CARD,
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
        add_shadow(card)
        text_box(s, sx + Cm(0.6), y + Cm(0.5), sw - Cm(1.2), Cm(2.2), num,
                 size=46, bold=True, color=TEAL, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, sx + Cm(0.7), y + Cm(2.9), sw - Cm(1.2), Cm(1.4), lbl,
                 size=18, color=WHITE, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.MIDDLE)


def slide_competition(prs):
    s = base_slide(prs, 6)
    headline(s,
             "We Don't Just Give You a Tool. We Give You a Home.", size=25)

    headers = ["Feature", "TribeLink", "Slack", "Notion", "Asana"]
    rows = [
        ("AI Team Matching", [1, 0, 0, 0]),
        ("Voting System", [1, 1, 0, 0]),
        ("Contribution Tracking", [1, 0, 1, 1]),
        ("Carbon Dashboard", [1, 0, 0, 0]),
        ("Virtual HQ", [1, 0, 0, 0]),
    ]

    nrows = len(rows) + 1
    ncols = len(headers)
    tx, ty = Cm(1.8), Cm(4.7)
    tw, th = Cm(30.3), Cm(12.5)
    tbl_shape = s.shapes.add_table(nrows, ncols, tx, ty, tw, th)
    table = tbl_shape.table

    # strip default table style for full manual control
    tblPr = table._tbl.tblPr
    tblPr.set("firstRow", "0")
    tblPr.set("bandRow", "0")
    sid = tblPr.find(qn("a:tableStyleId"))
    if sid is not None:
        tblPr.remove(sid)

    table.columns[0].width = Cm(9.9)
    rest = (Cm(30.3) - Cm(9.9))
    for c in range(1, ncols):
        table.columns[c].width = int(rest / (ncols - 1))
    table.rows[0].height = Cm(1.7)
    for r in range(1, nrows):
        table.rows[r].height = int((Cm(12.5) - Cm(1.7)) / (nrows - 1))

    # header row
    for c, head in enumerate(headers):
        is_tl = (head == "TribeLink")
        style_cell(table.cell(0, c), head, size=16, bold=True,
                   color=NAVY if is_tl else WHITE,
                   fill=TEAL if is_tl else RGBColor(0x12, 0x24, 0x36),
                   align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)

    for ri, (feat, marks) in enumerate(rows, start=1):
        band = CARD if ri % 2 else RGBColor(0x18, 0x27, 0x37)
        style_cell(table.cell(ri, 0), feat, size=14, bold=True, color=WHITE,
                   fill=band, align=PP_ALIGN.LEFT)
        for ci, ok in enumerate(marks, start=1):
            is_tl = (ci == 1)
            mark = "✔" if ok else "✘"
            if ok:
                col = TEAL if is_tl else GREEN_OK
            else:
                col = MUTED
            style_cell(table.cell(ri, ci), mark, size=20 if is_tl else 18,
                       bold=is_tl, color=col,
                       fill=RGBColor(0x10, 0x3A, 0x36) if is_tl else band)


def slide_business(prs):
    s = base_slide(prs, 7)
    headline(s, "Three Revenue Streams. Built to Scale.", size=30)

    pillars = [
        ("\U0001F4B3", "SaaS Subscriptions",
         [("Starter", "$12"), ("Pro", "$29"), ("Enterprise", "Custom")]),
        ("⭐", "Premium Features",
         [("Advanced AI", ""), ("Voting archives", ""),
          ("Remote support", "")]),
        ("\U0001F331", "Ethical Advertising",
         [("Eco-brand", ""), ("partnerships via", ""),
          ("Green Workspace", "")]),
    ]
    n = 3
    gap = Cm(0.9)
    left = Cm(1.6)
    total = SLIDE_W - left - Cm(1.2)
    cw = (total - gap * (n - 1)) / n
    top = Cm(5.0)
    ch = Cm(9.4)
    for i, (icon, title, items) in enumerate(pillars):
        x = left + i * (cw + gap)
        card = add_rect(s, x, top, cw, ch, fill=CARD,
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
        add_shadow(card)
        add_rect(s, x, top, cw, Cm(0.3), fill=TEAL)  # top accent
        text_box(s, x, top + Cm(0.8), cw, Cm(1.8), icon, size=44,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x + Cm(0.4), top + Cm(2.9), cw - Cm(0.8), Cm(1.2), title,
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        # detail lines
        paras = []
        for a, b in items:
            if b:
                paras.append([(f"{a}  ", {"color": GREY, "size": 15}),
                              (b, {"color": TEAL, "size": 16, "bold": True})])
            else:
                paras.append([(a, {"color": GREY, "size": 15})])
        text_box(s, x + Cm(0.7), top + Cm(4.6), cw - Cm(1.4), Cm(4.0), paras,
                 align=PP_ALIGN.CENTER, line_spacing=1.5)

    # flow arrow at bottom
    bar = add_rect(s, Cm(4.0), Cm(15.6), Cm(25.87), Cm(1.5), fill=TEAL,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    add_shadow(bar)
    text_box(s, Cm(4.0), Cm(15.6), Cm(25.87), Cm(1.5),
             "Freemium  →  Conversion  →  Upsell  →  Partnership",
             size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


def slide_traction(prs):
    s = base_slide(prs, 8)
    headline(s, "We Have Proof.", size=32)

    # left metric boxes
    metrics = [
        ("847", "Beta Participants  ·  22 Countries"),
        ("78%", "Retention Rate  (vs. 55% industry avg)"),
        ("+62", "NPS Score"),
    ]
    mx = Cm(1.6)
    mw = Cm(13.5)
    mh = Cm(2.8)
    my0 = Cm(5.2)
    for i, (num, lbl) in enumerate(metrics):
        y = my0 + i * (mh + Cm(0.65))
        card = add_rect(s, mx, y, mw, mh, fill=CARD,
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
        add_shadow(card)
        text_box(s, mx + Cm(0.5), y, Cm(4.3), mh, num, size=40, bold=True,
                 color=TEAL, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, mx + Cm(5.0), y, mw - Cm(5.4), mh, lbl, size=15,
                 color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.1)

    # right timeline
    items = [
        ("Q3 2026", "Global Public\nBeta"),
        ("Q4 2026", "Slack & Notion\nIntegration"),
        ("Q1 2027", "Virtual HQ\n2.0"),
        ("Q2 2027", "Series A\nReady"),
    ]
    tx0 = Cm(16.4)
    tw = Cm(16.0)
    line_y = Cm(9.0)
    add_rect(s, tx0, line_y, tw, Cm(0.08), fill=TEAL)
    step = tw / len(items)
    for i, (q, desc) in enumerate(items):
        cx = tx0 + step * i + step / 2
        node = add_rect(s, cx - Cm(0.35), line_y - Cm(0.27), Cm(0.7), Cm(0.7),
                        fill=TEAL, line=NAVY, line_w=Pt(2),
                        shape=MSO_SHAPE.OVAL)
        text_box(s, cx - step / 2, line_y - Cm(2.0), step, Cm(1.6), q,
                 size=15, bold=True, color=TEAL, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.BOTTOM)
        text_box(s, cx - step / 2, line_y + Cm(0.7), step, Cm(2.4), desc,
                 size=12.5, color=GREY, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.TOP, line_spacing=1.05)

    # pilot logos placeholder
    pilots = ["Fintech  SG", "EdTech  Berlin", "Agency  São Paulo"]
    pw = Cm(7.4)
    pgap = Cm(0.9)
    ptotal = pw * 3 + pgap * 2
    px0 = (SLIDE_W - ptotal) / 2
    py = Cm(15.8)
    text_box(s, px0, Cm(14.85), ptotal, Cm(0.8),
             "Pilot partners", size=12, color=GREY, italic=True,
             align=PP_ALIGN.CENTER)
    for i, name in enumerate(pilots):
        x = px0 + i * (pw + pgap)
        add_rect(s, x, py, pw, Cm(1.7), fill=None, line=GREY, line_w=Pt(1.25),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
        text_box(s, x, py, pw, Cm(1.7), name, size=14, bold=True, color=GREY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_financials(prs, img):
    s = base_slide(prs, 9)
    headline(s, "Raising $500K Seed Round", size=32)

    ih = Cm(11.0)
    s.shapes.add_picture(img, Cm(1.6), Cm(5.2), height=ih)

    # legend for donut
    legend = [("40%  Product Development", RGBColor(0x2D, 0x6C, 0xDF)),
              ("30%  Sales & Marketing", TEAL),
              ("20%  Operations & Talent", RGBColor(0x3C, 0xB3, 0x71)),
              ("10%  Legal & Compliance", RGBColor(0x90, 0xA4, 0xAE))]
    lx = Cm(1.6)
    ly = Cm(16.4)
    # arrange legend in two columns under chart
    for i, (txt, col) in enumerate(legend):
        cxx = lx + (i % 2) * Cm(6.6)
        cyy = ly + (i // 2) * Cm(0.85)
        add_rect(s, cxx, cyy + Cm(0.08), Cm(0.45), Cm(0.45), fill=col)
        text_box(s, cxx + Cm(0.65), cyy, Cm(5.8), Cm(0.7), txt, size=11,
                 color=GREY, anchor=MSO_ANCHOR.MIDDLE)

    # right projection table
    rows = [("Milestone", "Timeline"),
            ("10,000 paying users", "Month 12"),
            ("$2.4M ARR", "Month 18"),
            ("Series A target valuation", "$8–10M")]
    tx, ty = Cm(16.6), Cm(5.4)
    tw, th = Cm(15.7), Cm(6.5)
    tbl_shape = s.shapes.add_table(len(rows), 2, tx, ty, tw, th)
    table = tbl_shape.table
    tblPr = table._tbl.tblPr
    tblPr.set("firstRow", "0")
    tblPr.set("bandRow", "0")
    sid = tblPr.find(qn("a:tableStyleId"))
    if sid is not None:
        tblPr.remove(sid)
    table.columns[0].width = Cm(10.2)
    table.columns[1].width = Cm(5.5)
    for r in range(len(rows)):
        table.rows[r].height = int(Cm(6.5) / len(rows))
    for ri, (a, b) in enumerate(rows):
        if ri == 0:
            style_cell(table.cell(ri, 0), a, size=15, bold=True, color=NAVY,
                       fill=TEAL, align=PP_ALIGN.LEFT)
            style_cell(table.cell(ri, 1), b, size=15, bold=True, color=NAVY,
                       fill=TEAL, align=PP_ALIGN.CENTER)
        else:
            band = CARD if ri % 2 else RGBColor(0x18, 0x27, 0x37)
            style_cell(table.cell(ri, 0), a, size=14, color=WHITE, fill=band,
                       align=PP_ALIGN.LEFT)
            style_cell(table.cell(ri, 1), b, size=14, bold=True, color=TEAL,
                       fill=band, align=PP_ALIGN.CENTER)

    # callout box
    box = add_rect(s, Cm(16.6), Cm(13.2), Cm(15.7), Cm(2.6), fill=None,
                   line=TEAL, line_w=Pt(2), shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                   radius=0.12)
    text_box(s, Cm(16.6), Cm(13.2), Cm(15.7), Cm(2.6),
             "18-Month Path to Series A", size=22, bold=True, color=TEAL,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_closing(prs):
    s = base_slide(prs, 10)
    text_box(s, Cm(2.4), Cm(1.6), Cm(30), Cm(1.8),
             "Let's Build the Future of Work — Together.", size=32,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    # large gold quote
    text_box(s, Cm(3.0), Cm(5.0), Cm(27.87), Cm(2.6),
             "“Distance is never a barrier to great work.”",
             size=30, bold=True, italic=True, color=GOLD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Cm(14.43), Cm(8.0), Cm(5.0), Cm(0.08), fill=GOLD)

    # speaker sign-off blocks
    speakers = [("Speaker A", "CEO", "ceo@tribelink.io"),
                ("Speaker B", "CPO", "cpo@tribelink.io"),
                ("Speaker C", "CFO", "cfo@tribelink.io")]
    n = 3
    gap = Cm(0.9)
    left = Cm(2.5)
    total = SLIDE_W - 2 * left
    cw = (total - gap * (n - 1)) / n
    top = Cm(9.6)
    ch = Cm(3.4)
    for i, (name, role, email) in enumerate(speakers):
        x = left + i * (cw + gap)
        add_rect(s, x, top, cw, ch, fill=CARD,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
        add_rect(s, x + cw / 2 - Cm(0.9), top + Cm(0.4), Cm(1.8), Cm(0.08),
                 fill=GOLD)
        text_box(s, x, top + Cm(0.6), cw, Cm(1.0), name, size=19, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x, top + Cm(1.6), cw, Cm(0.8), role, size=14, bold=True,
                 color=GOLD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text_box(s, x, top + Cm(2.3), cw, Cm(0.9), email, size=13, color=GREY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # bottom: logo + website + QR placeholders
    by = Cm(15.6)
    logo = add_rect(s, Cm(2.6), by, Cm(2.4), Cm(2.4), fill=None, line=GOLD,
                    line_w=Pt(2), shape=MSO_SHAPE.OVAL)
    text_box(s, Cm(2.6), by, Cm(2.4), Cm(2.4), "TL", size=20, bold=True,
             color=GOLD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text_box(s, Cm(5.4), by, Cm(16), Cm(2.4),
             [[("TribeLink", {"color": WHITE, "size": 18, "bold": True})],
              [("www.tribelink.io", {"color": GREY, "size": 14})]],
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)

    qr = add_rect(s, Cm(28.9), by, Cm(2.4), Cm(2.4), fill=None, line=GREY,
                  line_w=Pt(1.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                  radius=0.1)
    text_box(s, Cm(28.9), by, Cm(2.4), Cm(2.4), "QR\nCODE", size=11,
             color=GREY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.0)


# --------------------------------------------------------------------------- #
def main():
    os.makedirs(ASSETS, exist_ok=True)
    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)

    overlap = os.path.join(ASSETS, "overlap.png")
    concentric = os.path.join(ASSETS, "concentric.png")
    donut = os.path.join(ASSETS, "donut.png")
    make_overlap_circles(overlap)
    make_concentric(concentric)
    make_donut(donut)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs, overlap)
    slide_features(prs)
    slide_market(prs, concentric)
    slide_competition(prs)
    slide_business(prs)
    slide_traction(prs)
    slide_financials(prs, donut)
    slide_closing(prs)

    prs.save(OUTPUT)
    print(f"Saved {OUTPUT} with {len(prs.slides._sldIdLst)} slides")


if __name__ == "__main__":
    main()
